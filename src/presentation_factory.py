from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def crear_presentacion():
    # Crear presentación
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Colores corporativos
    COLOR_PRIMARIO = RGBColor(30, 58, 138)  # Azul oscuro
    COLOR_ACENTO = RGBColor(16, 185, 129)   # Verde
    COLOR_ALERTA = RGBColor(239, 68, 68)    # Rojo
    COLOR_BLANCO = RGBColor(255, 255, 255)
    COLOR_GRIS = RGBColor(50, 50, 50)
    
    # ============================================
    # SLIDE 1: PORTADA
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARIO
    
    # Título principal
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(1)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "Optimización de Compras Internacionales"
    
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLANCO
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.6))
    tf = textbox.text_frame
    tf.text = "Pronóstico de Tipo de Cambio mediante Deep Learning"
    
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_ACENTO
    p.alignment = PP_ALIGN.CENTER
    
    # Autor y fecha
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    tf = textbox.text_frame
    tf.text = "[Tu Nombre] | Diciembre 2024 | Modelos no Lineales para Pronóstico"
    
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 2: EL PROBLEMA
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = textbox.text_frame
    tf.text = "El Problema: Incertidumbre Cambiaria en PyMEs"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARIO
    
    # Estadística impactante (caja izquierda)
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(4.5), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ALERTA
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.text = "70%\nde las PyMEs mexicanas importan insumos"
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_BLANCO
        paragraph.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1  # Middle
    
    # Problemas específicos
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(4.5), Inches(2))
    tf = textbox.text_frame
    tf.text = """📉 Volatilidad del tipo de cambio
💰 Erosión de márgenes de ganancia
❓ ¿Cuándo comprar divisas?
🌍 ¿Qué proveedor es más rentable?"""
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = COLOR_GRIS
        paragraph.space_after = Pt(12)
    
    # Impacto (caja derecha)
    shape = slide.shapes.add_shape(1, Inches(5.2), Inches(1.2), Inches(4.3), Inches(3.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 245, 255)
    shape.line.color.rgb = COLOR_PRIMARIO
    shape.line.width = Pt(2)
    
    textbox = slide.shapes.add_textbox(Inches(5.4), Inches(1.4), Inches(4), Inches(3.4))
    tf = textbox.text_frame
    tf.text = """Impacto Real:

- Volatilidad 2024: ±15% en 6 meses

- Una PyME que importa $1M USD puede perder hasta $150K MXN por mal timing

- Sin herramientas de pronóstico, las decisiones son reactivas, no estratégicas"""
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = COLOR_GRIS
        paragraph.space_after = Pt(10)
    
    # ============================================
    # SLIDE 3: LA SOLUCIÓN
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = textbox.text_frame
    tf.text = "La Solución: Sistema Predictivo Inteligente"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARIO
    
    # Columnas
    columnas = [
        {
            'left': 0.5,
            'color': RGBColor(240, 253, 244),
            'border': COLOR_ACENTO,
            'titulo': '📊 DATOS',
            'contenido': """✓ Banxico API
✓ Tipo de cambio
✓ TIIE, UDIS, CETES
✓ 2024-2025
✓ 350+ observaciones"""
        },
        {
            'left': 3.6,
            'color': RGBColor(240, 245, 255),
            'border': COLOR_PRIMARIO,
            'titulo': '🤖 MODELO',
            'contenido': """✓ LSTM 3 capas
✓ 60 días históricos
✓ Transformer (comparación)
✓ MAPE: X.X%
✓ Intervalos 50/80/95%"""
        },
        {
            'left': 6.7,
            'color': RGBColor(254, 242, 242),
            'border': COLOR_ALERTA,
            'titulo': '📈 ENTREGABLES',
            'contenido': """✓ Pronóstico 30 días
✓ 3 escenarios
✓ Dashboard interactivo
✓ Alertas automáticas"""
        }
    ]
    
    for col in columnas:
        # Caja
        shape = slide.shapes.add_shape(1, Inches(col['left']), Inches(1.2), Inches(2.8), Inches(3.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = col['color']
        shape.line.color.rgb = col['border']
        shape.line.width = Pt(2)
        
        # Contenido
        textbox = slide.shapes.add_textbox(Inches(col['left'] + 0.1), Inches(1.4), Inches(2.6), Inches(3.4))
        tf = textbox.text_frame
        tf.text = col['titulo'] + '\n\n' + col['contenido']
        
        for i, paragraph in enumerate(tf.paragraphs):
            if i == 0:
                paragraph.font.size = Pt(16)
                paragraph.font.bold = True
            else:
                paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = COLOR_GRIS
            paragraph.space_after = Pt(8)
    
    # ============================================
    # SLIDE 4: RESULTADOS
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = textbox.text_frame
    tf.text = "Resultados: Precisión Accionable"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARIO
    
    # Métricas
    metricas = [
        {'left': 0.5, 'valor': 'X.XX MXN', 'label': 'RMSE'},
        {'left': 3.6, 'valor': 'X.X%', 'label': 'MAPE'},
        {'left': 6.7, 'valor': '9X%', 'label': 'Precisión'}
    ]
    
    for m in metricas:
        # Caja métrica
        shape = slide.shapes.add_shape(1, Inches(m['left']), Inches(1.2), Inches(2.8), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_ACENTO
        shape.line.fill.background()
        
        # Valor
        textbox = slide.shapes.add_textbox(Inches(m['left']), Inches(1.3), Inches(2.8), Inches(0.5))
        tf = textbox.text_frame
        tf.text = m['valor']
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        textbox = slide.shapes.add_textbox(Inches(m['left']), Inches(1.9), Inches(2.8), Inches(0.4))
        tf = textbox.text_frame
        tf.text = m['label']
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.CENTER
    
    # Valor de negocio
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(2.6), Inches(9), Inches(2.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(249, 250, 251)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    shape.line.width = Pt(1)
    
    textbox = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(8.6), Inches(2.3))
    tf = textbox.text_frame
    tf.text = """💰 Valor de Negocio Tangible:

- Optimización de timing: Comprar divisas en ventanas óptimas (±30 días anticipación)
- Selección dinámica de proveedores: Elegir proveedor según pronóstico de su divisa
- Cobertura inteligente: Reducir costos de hedging innecesario (ahorrar 2-5% anual)
- Impacto estimado: $50K-$200K MXN anuales para PyME con $1M USD en importaciones"""
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = COLOR_GRIS
        paragraph.space_after = Pt(10)
    
    # ============================================
    # SLIDE 5: CASO DE USO
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = textbox.text_frame
    tf.text = "Caso de Uso: Estrategia de Compras Multi-Proveedor"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARIO
    
    # Escenario
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(4), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 245, 255)
    shape.line.color.rgb = COLOR_PRIMARIO
    shape.line.width = Pt(2)
    
    textbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(3.6), Inches(1))
    tf = textbox.text_frame
    tf.text = """Escenario:
PyME manufacturera con 3 proveedores:
🇺🇸 USA (USD) | 🇪🇺 Europa (EUR) | 🇨🇳 China (CNY)"""
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = COLOR_GRIS
    
    # Sin modelo
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(2.5), Inches(4), Inches(2.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(254, 242, 242)
    shape.line.color.rgb = COLOR_ALERTA
    shape.line.width = Pt(2)
    
    textbox = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(3.6), Inches(2.5))
    tf = textbox.text_frame
    tf.text = """❌ Sin Modelo:

- Decisión basada en precio actual
- Compra reactiva (cuando se necesita)
- Sin visibilidad de tendencias
- Pérdidas por volatilidad: ~8-12%"""
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = COLOR_GRIS
        paragraph.space_after = Pt(8)
    
    # Con modelo
    shape = slide.shapes.add_shape(1, Inches(5.5), Inches(1.1), Inches(4), Inches(4.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 253, 244)
    shape.line.color.rgb = COLOR_ACENTO
    shape.line.width = Pt(3)
    
    textbox = slide.shapes.add_textbox(Inches(5.7), Inches(1.2), Inches(3.6), Inches(3.8))
    tf = textbox.text_frame
    tf.text = """✅ Con Modelo Predictivo:

1. Pronóstico 30 días de USD/MXN

2. Si USD proyecta baja:
   → Priorizar proveedor USA
   → Anticipar compra de divisas

3. Si USD proyecta alza:
   → Evaluar proveedores EUR/CNY
   → Diferir compra si es posible

Ahorro estimado: 5-8% anual"""
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = COLOR_GRIS
        paragraph.space_after = Pt(6)
    
    # ============================================
    # SLIDE 6: PRÓXIMOS PASOS
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = textbox.text_frame
    tf.text = "Implementación y Roadmap"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARIO
    
    # Fases
    fases = [
        {
            'left': 0.5,
            'color': COLOR_ACENTO,
            'text_color': COLOR_BLANCO,
            'titulo': 'Fase 1: MVP',
            'tiempo': '2-4 semanas',
            'contenido': '• Dashboard básico\n• Pronóstico USD/MXN\n• Alertas email'
        },
        {
            'left': 3.6,
            'color': RGBColor(240, 245, 255),
            'text_color': COLOR_PRIMARIO,
            'titulo': 'Fase 2: Escala',
            'tiempo': '1-2 meses',
            'contenido': '• Multi-divisa (EUR, CNY)\n• API integración\n• Recomendaciones automáticas'
        },
        {
            'left': 6.7,
            'color': RGBColor(240, 245, 255),
            'text_color': COLOR_PRIMARIO,
            'titulo': 'Fase 3: IA Avanzada',
            'tiempo': '3-6 meses',
            'contenido': '• Sentiment analysis\n• Eventos geopolíticos\n• Optimización de portafolio'
        }
    ]
    
    for fase in fases:
        # Caja fase
        shape = slide.shapes.add_shape(1, Inches(fase['left']), Inches(1.2), Inches(2.8), Inches(3.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fase['color']
        shape.line.color.rgb = COLOR_PRIMARIO
        shape.line.width = Pt(2)
        
        # Título
        textbox = slide.shapes.add_textbox(Inches(fase['left'] + 0.1), Inches(1.3), Inches(2.6), Inches(0.4))
        tf = textbox.text_frame
        tf.text = fase['titulo']
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = fase['text_color']
        p.alignment = PP_ALIGN.CENTER
        
        # Tiempo
        textbox = slide.shapes.add_textbox(Inches(fase['left'] + 0.1), Inches(1.8), Inches(2.6), Inches(0.3))
        tf = textbox.text_frame
        tf.text = fase['tiempo']
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = fase['text_color'] if fase['text_color'] == COLOR_BLANCO else RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
        
        # Contenido
        textbox = slide.shapes.add_textbox(Inches(fase['left'] + 0.2), Inches(2.3), Inches(2.4), Inches(2.3))
        tf = textbox.text_frame
        tf.text = fase['contenido']
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = fase['text_color'] if fase['text_color'] == COLOR_BLANCO else COLOR_GRIS
            paragraph.space_after = Pt(6)
    
    # Call to action
    shape = slide.shapes.add_shape(1, Inches(2), Inches(5), Inches(6), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ALERTA
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.text = "🚀 Comencemos con Fase 1 - ROI esperado en 3-6 meses"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLANCO
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1
    
    # ============================================
    # Guardar
    # ============================================
    import os
    # Obtener el directorio donde está el script
    script_dir = os.path.dirname(os.path.abspath(__file__))
    filename = os.path.join(script_dir, 'Presentacion_Pronostico_Financiero.pptx')
    prs.save(filename)
    print(f"✅ Presentación creada exitosamente en: {filename}")

if __name__ == "__main__":
    crear_presentacion()